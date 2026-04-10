"""Tests for the PPTX builder module."""

from pathlib import Path

from PIL import Image

from pptx import Presentation
from pptx.util import Pt

from png2pptx.models import SlideData, TextBlock, WordBox
from png2pptx.pptx_builder import (
    SLIDE_LONG_EDGE_EMU,
    _compute_slide_dimensions,
    _fit_font_size_px,
    build_pptx,
)


def _make_test_png(tmp_path: Path) -> Path:
    """Create a minimal test PNG."""
    img = Image.new("RGB", (800, 600), color="white")
    path = tmp_path / "test_slide.png"
    img.save(path)
    return path


def test_creates_file(tmp_path):
    """build_pptx should create a .pptx file."""
    img_path = _make_test_png(tmp_path)
    out = tmp_path / "out.pptx"
    result = build_pptx(
        [SlideData(image_path=img_path, image_width=800, image_height=600)],
        out,
    )
    assert result.exists()
    assert result.suffix == ".pptx"


def test_slide_count(tmp_path):
    """Should create one slide per input image."""
    img_path = _make_test_png(tmp_path)
    out = tmp_path / "out.pptx"
    slides = [
        SlideData(image_path=img_path, image_width=800, image_height=600),
        SlideData(image_path=img_path, image_width=800, image_height=600),
    ]
    build_pptx(slides, out)
    prs = Presentation(str(out))
    assert len(prs.slides) == 2


def test_text_boxes_created(tmp_path):
    """Should create text boxes for text blocks."""
    img_path = _make_test_png(tmp_path)
    out = tmp_path / "out.pptx"

    block = TextBlock(
        words=[
            WordBox(text="Hello", x=100, y=100, width=200, height=40,
                    confidence=95.0),
        ],
        color=(255, 0, 0),
    )
    slide_data = SlideData(
        image_path=img_path,
        image_width=800,
        image_height=600,
        text_blocks=[block],
    )
    build_pptx([slide_data], out)

    prs = Presentation(str(out))
    slide = prs.slides[0]
    # Should have background image + text box = at least 2 shapes
    assert len(slide.shapes) >= 2

    # Find the text box (not the picture)
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    assert len(text_shapes) == 1
    assert "Hello" in text_shapes[0].text_frame.text


def test_empty_slides(tmp_path):
    """Should handle empty slides list."""
    out = tmp_path / "empty.pptx"
    build_pptx([], out)
    prs = Presentation(str(out))
    assert len(prs.slides) == 0


def test_portrait_slide_uses_long_edge_for_height():
    """Portrait images should get a tall slide, not a squished one."""
    slide_w, slide_h = _compute_slide_dimensions(1824, 2336)
    assert slide_h == SLIDE_LONG_EDGE_EMU
    assert slide_w < slide_h  # Portrait: width < height


def test_landscape_slide_uses_long_edge_for_width():
    """Landscape images should get a wide slide."""
    slide_w, slide_h = _compute_slide_dimensions(1920, 1080)
    assert slide_w == SLIDE_LONG_EDGE_EMU
    assert slide_w > slide_h


def test_font_size_scales_with_image(tmp_path):
    """Font size should be proportional to word height on the slide."""
    img_path = _make_test_png(tmp_path)
    out = tmp_path / "font_test.pptx"

    # Word that is 40px tall in an 600px-tall image
    block = TextBlock(
        words=[
            WordBox(text="Test", x=50, y=50, width=150, height=40,
                    confidence=95.0),
        ],
    )
    slide_data = SlideData(
        image_path=img_path, image_width=800, image_height=600,
        text_blocks=[block],
    )
    build_pptx([slide_data], out)

    prs = Presentation(str(out))
    text_shapes = [s for s in prs.slides[0].shapes if s.has_text_frame]
    font_size = text_shapes[0].text_frame.paragraphs[0].runs[0].font.size

    # 40px / 600px image on a ~13.33" slide ≈ 0.889 inches of box height
    # 0.889" * 72pt/in = 64pt, * 0.85 = 54.4pt
    # Should be in a reasonable range (not 6pt or 200pt)
    assert Pt(20) < font_size < Pt(96)


def test_width_aware_font_fitting_shrinks_long_text():
    """Longer text should fit at a smaller size when the box width is fixed."""
    short_block = TextBlock(
        words=[
            WordBox(text="Hi", x=10, y=10, width=90, height=36, confidence=95.0),
        ]
    )
    long_block = TextBlock(
        words=[
            WordBox(
                text="SignificantlyLongerLabel",
                x=10,
                y=10,
                width=90,
                height=36,
                confidence=95.0,
            ),
        ]
    )

    short_size = _fit_font_size_px(short_block.text, short_block)
    long_size = _fit_font_size_px(long_block.text, long_block)

    assert long_size < short_size
