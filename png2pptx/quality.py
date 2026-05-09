from __future__ import annotations

import json
from dataclasses import asdict, dataclass, field
from pathlib import Path

import numpy as np
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

from .inpaint import remove_text
from .layout import group_into_blocks
from .models import SlideData
from .ocr import extract_words
from .pptx_builder import build_pptx
from .styles import extract_text_colors

_GENERATED_EXAMPLE_SUFFIXES = (
    "_baseline_clean",
    "_current_clean",
    "_compare",
    "_overlay",
    "_review",
)

_REVIEW_LABEL_HEIGHT = 22
_REVIEW_PADDING = 8
_OVERLAY_BOX_COLOR = (255, 48, 48)
_OVERLAY_LABEL_COLOR = (255, 235, 59)


@dataclass
class QualityResult:
    name: str
    image_size: tuple[int, int]
    ocr_mode: str
    confidence: float
    word_count: int
    block_count: int
    source_changed_ratio: float
    source_mae: float
    pptx_text_shapes: int
    pptx_text_runs: int
    sample_text: list[str] = field(default_factory=list)
    clean_path: str = ""
    pptx_path: str = ""
    overlay_path: str = ""
    review_path: str = ""
    baseline_path: str | None = None
    baseline_mae: float | None = None
    baseline_max_diff: int | None = None
    baseline_changed_ratio: float | None = None


def discover_example_images(examples_dir: str | Path) -> list[Path]:
    directory = Path(examples_dir)
    candidates: list[Path] = []
    for path in sorted(directory.glob("*.png")):
        if any(path.stem.endswith(suffix) for suffix in _GENERATED_EXAMPLE_SUFFIXES):
            continue
        candidates.append(path)
    return candidates


def run_quality_loop(
    examples_dir: str | Path,
    output_dir: str | Path,
    confidence: float = 40.0,
    lang: str = "eng",
    ocr_mode: str = "aggressive",
    remove_background_text: bool = True,
) -> list[QualityResult]:
    examples_dir = Path(examples_dir)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    results: list[QualityResult] = []
    for image_path in discover_example_images(examples_dir):
        results.append(
            _process_example(
                image_path=image_path,
                output_dir=output_dir,
                confidence=confidence,
                lang=lang,
                ocr_mode=ocr_mode,
                remove_background_text=remove_background_text,
            )
        )

    summary_path = output_dir / "summary.json"
    summary_path.write_text(
        json.dumps([asdict(result) for result in results], indent=2),
        encoding="utf-8",
    )
    return results


def _process_example(
    image_path: Path,
    output_dir: Path,
    confidence: float,
    lang: str,
    ocr_mode: str,
    remove_background_text: bool,
) -> QualityResult:
    words, img_width, img_height = extract_words(
        image_path,
        confidence_threshold=confidence,
        lang=lang,
        ocr_mode=ocr_mode,
    )
    blocks = group_into_blocks(words)
    blocks = extract_text_colors(image_path, blocks)

    clean_path = output_dir / f"{image_path.stem}_current_clean.png"
    background_path = image_path
    if remove_background_text and blocks:
        background_path = remove_text(image_path, blocks, output_path=clean_path)
    else:
        clean_path = image_path

    pptx_path = output_dir / f"{image_path.stem}.pptx"
    build_pptx(
        [
            SlideData(
                image_path=background_path,
                image_width=img_width,
                image_height=img_height,
                text_blocks=blocks,
            )
        ],
        pptx_path,
    )

    overlay_path = output_dir / f"{image_path.stem}_overlay.png"
    render_pptx_overlay(pptx_path, image_path, overlay_path)

    baseline_path = image_path.with_name(f"{image_path.stem}_baseline_clean.png")
    review_path = output_dir / f"{image_path.stem}_review.png"
    build_review_image(
        image_path=image_path,
        clean_path=clean_path,
        overlay_path=overlay_path,
        output_path=review_path,
        baseline_path=baseline_path if baseline_path.exists() else None,
    )

    source_metrics = compute_image_delta_metrics(image_path, clean_path)
    baseline_metrics: dict[str, float | int] | None = None
    if baseline_path.exists():
        baseline_metrics = compute_image_delta_metrics(baseline_path, clean_path)

    pptx_summary = summarize_pptx_text(pptx_path)
    return QualityResult(
        name=image_path.stem,
        image_size=(img_width, img_height),
        ocr_mode=ocr_mode,
        confidence=confidence,
        word_count=len(words),
        block_count=len(blocks),
        source_changed_ratio=source_metrics["changed_ratio"],
        source_mae=source_metrics["mae"],
        pptx_text_shapes=pptx_summary["text_shapes"],
        pptx_text_runs=pptx_summary["text_runs"],
        sample_text=pptx_summary["sample_text"],
        clean_path=str(clean_path),
        pptx_path=str(pptx_path),
        overlay_path=str(overlay_path),
        review_path=str(review_path),
        baseline_path=str(baseline_path) if baseline_path.exists() else None,
        baseline_mae=None if baseline_metrics is None else float(baseline_metrics["mae"]),
        baseline_max_diff=None if baseline_metrics is None else int(baseline_metrics["max_diff"]),
        baseline_changed_ratio=None
        if baseline_metrics is None
        else float(baseline_metrics["changed_ratio"]),
    )


def compute_image_delta_metrics(
    reference_image: str | Path,
    candidate_image: str | Path,
) -> dict[str, float | int]:
    reference = np.asarray(Image.open(reference_image).convert("RGB"), dtype=np.int16)
    candidate = np.asarray(Image.open(candidate_image).convert("RGB"), dtype=np.int16)
    if reference.shape != candidate.shape:
        raise ValueError("Images must have the same dimensions to compute delta metrics.")

    diff = np.abs(reference - candidate)
    changed_ratio = float(np.any(diff != 0, axis=2).mean())
    mae = float(diff.mean())
    max_diff = int(diff.max())
    return {
        "changed_ratio": changed_ratio,
        "mae": mae,
        "max_diff": max_diff,
    }


def summarize_pptx_text(pptx_path: str | Path) -> dict[str, int | list[str]]:
    presentation = Presentation(str(pptx_path))
    text_shapes = 0
    text_runs = 0
    sample_text: list[str] = []

    for shape in presentation.slides[0].shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        text_shapes += 1
        sample_text.append(text)
        for paragraph in shape.text_frame.paragraphs:
            text_runs += len(paragraph.runs)

    return {
        "text_shapes": text_shapes,
        "text_runs": text_runs,
        "sample_text": sample_text[:8],
    }


def render_pptx_overlay(
    pptx_path: str | Path,
    image_path: str | Path,
    output_path: str | Path,
) -> Path:
    presentation = Presentation(str(pptx_path))
    slide = presentation.slides[0]

    image = Image.open(image_path).convert("RGB")
    image_width, image_height = image.size
    draw = ImageDraw.Draw(image)
    font = ImageFont.load_default()

    scale_x = image_width / presentation.slide_width
    scale_y = image_height / presentation.slide_height

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue

        left = int(shape.left * scale_x)
        top = int(shape.top * scale_y)
        right = int((shape.left + shape.width) * scale_x)
        bottom = int((shape.top + shape.height) * scale_y)

        draw.rectangle([left, top, right, bottom], outline=_OVERLAY_BOX_COLOR, width=2)

        font_size_pt = _first_run_font_size_pt(shape)
        label = f"{font_size_pt:.0f}pt {text[:36]}".rstrip()
        draw.text(
            (left + 2, max(0, top - 14)),
            label,
            fill=_OVERLAY_LABEL_COLOR,
            font=font,
        )

    output_path = Path(output_path)
    image.save(output_path)
    return output_path


def build_review_image(
    image_path: str | Path,
    clean_path: str | Path,
    overlay_path: str | Path,
    output_path: str | Path,
    baseline_path: str | Path | None = None,
) -> Path:
    panels = [
        ("source", Image.open(image_path).convert("RGB")),
        ("current_clean", Image.open(clean_path).convert("RGB")),
        ("overlay", Image.open(overlay_path).convert("RGB")),
    ]
    if baseline_path is not None:
        panels.insert(1, ("baseline_clean", Image.open(baseline_path).convert("RGB")))

    output = _compose_review_panels(panels)
    output_path = Path(output_path)
    output.save(output_path)
    return output_path


def _compose_review_panels(panels: list[tuple[str, Image.Image]]) -> Image.Image:
    max_width = max(image.width for _, image in panels)
    max_height = max(image.height for _, image in panels)
    columns = 2 if len(panels) > 2 else len(panels)
    rows = (len(panels) + columns - 1) // columns

    sheet_width = columns * (max_width + (_REVIEW_PADDING * 2))
    sheet_height = rows * (max_height + _REVIEW_LABEL_HEIGHT + (_REVIEW_PADDING * 2))
    sheet = Image.new("RGB", (sheet_width, sheet_height), color=(24, 24, 24))

    font = ImageFont.load_default()
    draw = ImageDraw.Draw(sheet)
    for index, (label, image) in enumerate(panels):
        row = index // columns
        column = index % columns
        x = column * (max_width + (_REVIEW_PADDING * 2)) + _REVIEW_PADDING
        y = row * (max_height + _REVIEW_LABEL_HEIGHT + (_REVIEW_PADDING * 2)) + _REVIEW_PADDING
        draw.text((x, y), label, fill=(240, 240, 240), font=font)
        sheet.paste(image, (x, y + _REVIEW_LABEL_HEIGHT))
        draw.rectangle(
            [x - 1, y + _REVIEW_LABEL_HEIGHT - 1, x + image.width, y + _REVIEW_LABEL_HEIGHT + image.height],
            outline=(96, 96, 96),
            width=1,
        )

    return sheet


def _first_run_font_size_pt(shape) -> float:
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                return float(run.font.size / 12700)
    return 0.0
