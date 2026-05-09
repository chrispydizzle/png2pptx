"""Build a PPTX file with image backgrounds and editable text overlays."""

from __future__ import annotations

from functools import lru_cache
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt
from PIL import ImageFont

import re

from .models import SlideData, TextBlock

# Pattern to clean OCR artifacts from text edges
_EDGE_JUNK_RE = re.compile(r'^[_\-—–\s]+|[_\-—–\s]+$')
_FONT_FAMILY = "Calibri"
_FONT_WIDTH_PADDING = 1.05
_FONT_HEIGHT_PADDING = 1.02
_TEXTBOX_WIDTH_PADDING_RATIO = 0.05
_TEXTBOX_HEIGHT_PADDING_RATIO = 0.12

_FONT_PATH_CANDIDATES = (
    r"C:\Windows\Fonts\calibri.ttf",
    r"C:\Windows\Fonts\arial.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    "/Library/Fonts/Arial.ttf",
)


EMU_PER_INCH = 914400
EMU_PER_PT = 12700

# Standard 16:9 slide long edge in EMU (≈ 13.33 inches).
# We use this for the LONG edge of the slide so that portrait images
# get a generously-sized slide instead of being squashed.
SLIDE_LONG_EDGE_EMU = 12192000  # 13.33 in


def _compute_slide_dimensions(
    img_width: int, img_height: int
) -> tuple[int, int]:
    """Return (slide_width, slide_height) in EMU matching the image aspect ratio.

    The longer image dimension maps to SLIDE_LONG_EDGE_EMU; the shorter
    dimension is computed from the aspect ratio.
    """
    if img_width >= img_height:
        # Landscape or square
        slide_w = SLIDE_LONG_EDGE_EMU
        slide_h = int(SLIDE_LONG_EDGE_EMU * img_height / img_width)
    else:
        # Portrait
        slide_h = SLIDE_LONG_EDGE_EMU
        slide_w = int(SLIDE_LONG_EDGE_EMU * img_width / img_height)

    return slide_w, slide_h


def build_pptx(slides: list[SlideData], output_path: str | Path) -> Path:
    """Build a PPTX file from slide data.

    Args:
        slides: List of SlideData, each becomes one slide.
        output_path: Where to save the .pptx file.

    Returns:
        Path to the created file.
    """
    output_path = Path(output_path)
    prs = Presentation()

    if slides:
        first = slides[0]
        slide_w, slide_h = _compute_slide_dimensions(
            first.image_width, first.image_height
        )
        prs.slide_width = slide_w
        prs.slide_height = slide_h

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for slide_data in slides:
        slide = prs.slides.add_slide(blank_layout)
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # Add the PNG as a full-bleed background image
        slide.shapes.add_picture(
            str(slide_data.image_path),
            left=0,
            top=0,
            width=slide_w,
            height=slide_h,
        )

        # Scale factors: image pixels → slide EMUs
        scale_x = slide_w / slide_data.image_width
        scale_y = slide_h / slide_data.image_height

        # Add text boxes for each block/group
        for block_group in _group_blocks_for_rendering(slide_data.text_blocks):
            _add_text_group(slide, block_group, scale_x, scale_y)

    prs.save(str(output_path))
    return output_path


def _add_text_group(
    slide,
    blocks: list[TextBlock],
    scale_x: float,
    scale_y: float,
) -> None:
    if len(blocks) == 1:
        _add_text_block(slide, blocks[0], scale_x, scale_y)
        return

    left_px = min(block.x for block in blocks)
    top_px = min(block.y for block in blocks)
    right_px = max(block.right for block in blocks)
    bottom_px = max(block.bottom for block in blocks)

    left_emu = int(left_px * scale_x)
    top_emu = int(top_px * scale_y)
    width_emu = int((right_px - left_px) * scale_x)
    height_emu = int((bottom_px - top_px) * scale_y)

    h_pad = int(width_emu * _TEXTBOX_WIDTH_PADDING_RATIO)
    width_emu += h_pad
    v_pad = max(Emu(1), int(height_emu * _TEXTBOX_HEIGHT_PADDING_RATIO))
    top_emu = max(0, top_emu - (v_pad // 2))
    height_emu += v_pad

    txbox = slide.shapes.add_textbox(
        Emu(left_emu), Emu(top_emu), Emu(width_emu), Emu(height_emu)
    )

    # Transparent background
    txbox.fill.background()

    tf = txbox.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.auto_size = None
    _configure_paragraph_defaults(tf.paragraphs[0])

    for index, block in enumerate(sorted(blocks, key=lambda item: (item.y, item.x))):
        para = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        _configure_paragraph_defaults(para)
        para.alignment = PP_ALIGN.LEFT

        clean_text = _EDGE_JUNK_RE.sub('', block.text)
        rendered_text = clean_text if clean_text else block.text

        run = para.add_run()
        run.text = rendered_text

        fitted_font_px = _fit_font_size_px(rendered_text, block)
        font_size_pt = (fitted_font_px * scale_y) / EMU_PER_PT
        font_size_pt = max(6.0, min(font_size_pt, 96.0))
        run.font.size = Pt(font_size_pt)
        run.font.name = _FONT_FAMILY

        r, g, b = block.color
        run.font.color.rgb = RGBColor(r, g, b)

        if index + 1 < len(blocks):
            next_block = blocks[index + 1]
            gap_px = max(0, next_block.y - block.bottom)
            para.space_after = Pt((gap_px * scale_y) / EMU_PER_PT)


def _add_text_block(
    slide,
    block: TextBlock,
    scale_x: float,
    scale_y: float,
) -> None:
    """Add a transparent, single-line text box to the slide for one TextBlock."""
    clean_text = _EDGE_JUNK_RE.sub('', block.text)
    rendered_text = clean_text if clean_text else block.text

    # Convert block pixel coords → slide EMU coords
    left_emu = int(block.x * scale_x)
    top_emu = int(block.y * scale_y)
    width_emu = int(block.width * scale_x)
    height_emu = int(block.height * scale_y)

    # Minimal padding — just enough to avoid clipping the last character
    h_pad = int(width_emu * _TEXTBOX_WIDTH_PADDING_RATIO)
    width_emu += h_pad
    v_pad = max(Emu(1), int(height_emu * _TEXTBOX_HEIGHT_PADDING_RATIO))
    top_emu = max(0, top_emu - (v_pad // 2))
    height_emu += v_pad

    txbox = slide.shapes.add_textbox(
        Emu(left_emu), Emu(top_emu), Emu(width_emu), Emu(height_emu)
    )

    # Transparent background
    txbox.fill.background()

    tf = txbox.text_frame
    tf.word_wrap = False  # Single-line — never reflow
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    # Anchor text to the vertical middle of the box
    tf.auto_size = None
    _configure_paragraph_defaults(txbox.text_frame.paragraphs[0])

    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT

    run = para.add_run()
    run.text = rendered_text

    # Fit the font using both width and height so long labels shrink enough
    # while short headings can remain visually large.
    fitted_font_px = _fit_font_size_px(rendered_text, block)
    font_size_pt = (fitted_font_px * scale_y) / EMU_PER_PT
    font_size_pt = max(6.0, min(font_size_pt, 96.0))
    run.font.size = Pt(font_size_pt)
    run.font.name = _FONT_FAMILY

    # Font color
    r, g, b = block.color
    run.font.color.rgb = RGBColor(r, g, b)


def _configure_paragraph_defaults(paragraph) -> None:
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)


def _group_blocks_for_rendering(blocks: list[TextBlock]) -> list[list[TextBlock]]:
    if not blocks:
        return []

    sorted_blocks = sorted(blocks, key=lambda block: (block.y, block.x))
    groups: list[list[TextBlock]] = [[sorted_blocks[0]]]
    for block in sorted_blocks[1:]:
        current_group = groups[-1]
        if _should_group_blocks(current_group[-1], block):
            current_group.append(block)
        else:
            groups.append([block])

    return groups


def _should_group_blocks(previous: TextBlock, current: TextBlock) -> bool:
    vertical_gap = current.y - previous.bottom
    if vertical_gap < -2:
        return False
    if vertical_gap > max(previous.height, current.height) * 1.2 + 8.0:
        return False

    prev_center_x = previous.x + (previous.width / 2)
    current_center_x = current.x + (current.width / 2)
    left_close = abs(previous.x - current.x) <= max(28, max(previous.height, current.height) * 2.0)
    center_close = abs(prev_center_x - current_center_x) <= max(40, max(previous.width, current.width) * 0.12)
    if not left_close and not center_close:
        return False

    prev_font = max(previous.estimated_font_size_px, 1.0)
    current_font = max(current.estimated_font_size_px, 1.0)
    font_ratio = min(prev_font, current_font) / max(prev_font, current_font)
    if font_ratio < 0.8:
        return False

    if _color_distance(previous.color, current.color) > 60.0:
        return False

    return _has_multiline_continuation_signal(previous.text, current.text)


def _has_multiline_continuation_signal(previous_text: str, current_text: str) -> bool:
    previous_text = previous_text.strip()
    current_text = current_text.strip()
    if not previous_text or not current_text:
        return False

    current_first = current_text[:1]
    if current_first.islower() or current_first.isdigit():
        return True

    if previous_text.endswith((",", ";", ":", "-", "/", "(")):
        return True

    if len(previous_text) >= 28 and previous_text[-1:] not in ".!?":
        return True

    return False


def _color_distance(a: tuple[int, int, int], b: tuple[int, int, int]) -> float:
    return ((a[0] - b[0]) ** 2 + (a[1] - b[1]) ** 2 + (a[2] - b[2]) ** 2) ** 0.5


def _fit_font_size_px(text: str, block: TextBlock) -> float:
    """Find the largest rendered font size that fits inside the OCR box."""
    if not text.strip():
        return block.estimated_font_size_px

    target_width = max(1, int(block.width * _FONT_WIDTH_PADDING))
    target_height = max(1, int(block.height * _FONT_HEIGHT_PADDING))

    # Use the current geometry-based estimate as a seed, but search a wider range.
    seed = max(4, int(round(block.estimated_font_size_px)))
    low = 4
    high = max(seed * 2, int(block.height * 2.5), 12)
    best = low

    while low <= high:
        mid = (low + high) // 2
        measured_width, measured_height = _measure_text(text, mid)
        if measured_width <= target_width and measured_height <= target_height:
            best = mid
            low = mid + 1
        else:
            high = mid - 1

    return float(best)


@lru_cache(maxsize=1)
def _measurement_font_path() -> str | None:
    for candidate in _FONT_PATH_CANDIDATES:
        if Path(candidate).exists():
            return candidate
    return None


@lru_cache(maxsize=512)
def _measure_text(text: str, font_px: int) -> tuple[int, int]:
    """Measure single-line text in source-image pixels."""
    font_path = _measurement_font_path()
    if font_path is None:
        width = max(1, int(len(text) * font_px * 0.6))
        height = max(1, int(font_px))
        return width, height

    font = ImageFont.truetype(font_path, max(1, font_px))
    left, top, right, bottom = font.getbbox(text or "Ag")
    width = max(1, right - left)
    height = max(1, bottom - top)
    return width, height
